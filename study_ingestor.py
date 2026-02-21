import os
import sys
import json
from enum import Enum
from typing import List, Optional
from pydantic import BaseModel, Field, ValidationError
from pptx import Presentation
from pypdf import PdfReader

class Severity(str, Enum):
    low = "low"
    medium = "medium"
    high = "high"
    critical = "critical"

class Interaction(BaseModel):
    interacting_drug: str
    type: str
    severity: Severity
    management: str

class Drug(BaseModel):
    generic_name: str
    brand_name: str
    drug_class: str
    mechanism_of_action: str
    indications: str
    contraindications: str
    adverse_effects: str
    patient_education: str
    interactions: List[Interaction] = []

class Asset(BaseModel):
    condition: Optional[str] = None
    environment: Optional[str] = None
    action: Optional[str] = None
    drug: Optional[str] = None

class Scenario(BaseModel):
    id: str
    type: str # Relaxed validation to support 'mission' and 'extended_sata'
    text: Optional[str] = None
    question: Optional[str] = None
    options: Optional[List[str]] = None
    columns: Optional[List[str]] = None
    rows: Optional[List[dict]] = None
    actions: Optional[List[str]] = None
    conditions: Optional[List[str]] = None
    parameters: Optional[List[str]] = None
    correct_indices: Optional[List[int]] = None
    correct_col_idx: Optional[List[int]] = None # For matrix rows
    correct_actions: Optional[List[int]] = None
    correct_condition: Optional[List[int]] = None
    correct_parameters: Optional[List[int]] = None
    rationales: Optional[List[str]] = None
    assets: Optional[Asset] = None
    answer_type: Optional[str] = None
    
    # Mission Specific Fields
    numerical_id: Optional[int] = None
    title: Optional[str] = None
    category: Optional[str] = None
    difficulty: Optional[str] = None
    patient_profile: Optional[dict] = None
    exhibits: Optional[List[dict]] = None
    steps: Optional[List[dict]] = None
    prompt: Optional[str] = None # Used in mission steps instead of 'text'
    answer: Optional[object] = None # Used in mission steps instead of indices sometimes
    ragionale: Optional[str] = None # Typo handling or unified rationale string

class StudyIngestor:
    def __init__(self):
        self.accumulated_text = ""
        self.files_processed = []
        self.drugs: List[Drug] = []
        self.scenarios: List[Scenario] = []
        self.raw_data_dir = "raw_data"
        os.makedirs(self.raw_data_dir, exist_ok=True)

    def extract_text(self, file_path):
        if not os.path.exists(file_path):
            return f"Error: File '{file_path}' not found."
        
        ext = os.path.splitext(file_path)[1].lower()
        text = ""

        try:
            if ext == '.txt' or ext == '.md':
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
            elif ext == '.pptx':
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif ext == '.pdf':
                reader = PdfReader(file_path)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            else:
                return f"Error: Unsupported file type '{ext}'."
            
            # Save raw text
            filename = os.path.basename(file_path)
            raw_path = os.path.join(self.raw_data_dir, f"{filename}.txt")
            with open(raw_path, "w") as f:
                f.write(text)

            self.accumulated_text += f"\n--- SOURCE: {filename} ---\n" + text
            self.files_processed.append(filename)
            return f"Success: Read {len(text)} characters from {filename}. Saved to {raw_path}"
        
        except Exception as e:
            return f"Error processing {file_path}: {str(e)}"

    def validate_drug(self, drug_data: dict):
        try:
            drug = Drug(**drug_data)
            self.drugs.append(drug)
            return drug
        except ValidationError as e:
            print(f"Drug Validation Error: {e.json()}")
            return None

    def validate_scenario(self, scenario_data: dict):
        try:
            # Handle Matrix row/col structure adjustments if needed, 
            # but usually the LLM should output correct JSON structure matches.
            scenario = Scenario(**scenario_data)
            self.scenarios.append(scenario)
            return scenario
        except ValidationError as e:
            print(f"Scenario Validation Error: {e.json()}")
            return None

    def get_full_content(self):
        return self.accumulated_text

if __name__ == "__main__":
    ingestor = StudyIngestor()
    # Example testing block
    sample_drug = {
        "generic_name": "Warfarin",
        "brand_name": "Coumadin",
        "drug_class": "Anticoagulant",
        "mechanism_of_action": "Inhibits Vitamin K synthesis",
        "indications": "VTE prophylaxis",
        "contraindications": "Active bleeding",
        "adverse_effects": "Bleeding",
        "patient_education": "Monitor INR",
        "interactions": [
            {
                "interacting_drug": "Aspirin",
                "type": "Pharmacodynamic",
                "severity": "high",
                "management": "Monitor for bleeding signs"
            }
        ]
    }
    validated = ingestor.validate_drug(sample_drug)
    if validated:
        print(f"Successfully validated {validated.generic_name}")